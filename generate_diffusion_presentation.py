from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).with_name("presentation_reseaux_diffusion.pptx")

BG = RGBColor(246, 248, 251)
NAVY = RGBColor(17, 34, 64)
BLUE = RGBColor(42, 98, 176)
TEAL = RGBColor(27, 136, 145)
MINT = RGBColor(208, 240, 234)
ORANGE = RGBColor(231, 122, 72)
GOLD = RGBColor(231, 183, 71)
ROSE = RGBColor(239, 213, 221)
TEXT = RGBColor(49, 60, 74)
MUTED = RGBColor(104, 112, 124)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 223, 233)
GRAY_2 = RGBColor(196, 203, 214)
GRAY_3 = RGBColor(148, 156, 170)
GRAY_4 = RGBColor(104, 112, 124)


def set_run_font(run, name="Aptos", size=18, bold=False, color=TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.333),
        Inches(0.42),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = NAVY
    top_band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(10.6),
        0,
        Inches(2.733),
        Inches(0.42),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(12), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, name="Aptos Display", size=27, bold=True, color=NAVY)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.35), Inches(12), Inches(0.34))
        tf = sub.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        set_run_font(run, size=12, color=MUTED)


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12), Inches(0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    set_run_font(run, size=8, color=MUTED)


def add_text_block(slide, text, left, top, width, height, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, name="Aptos"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, name=name, size=size, bold=bold, color=color)
    return box


def add_panel(slide, left, top, width, height, title=None, fill_color=WHITE, line_color=LINE):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color
    panel.line.color.rgb = line_color
    panel.line.width = Pt(1.2)
    if title:
        add_text_block(slide, title, left + 0.18, top + 0.12, width - 0.36, 0.3, size=15, color=TEAL, bold=True, name="Aptos Display")
    return panel


def add_badge(slide, text, left, top, width, fill_color=MINT, text_color=TEAL):
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.38),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill_color
    badge.line.fill.background()
    add_text_block(slide, text, left + 0.12, top + 0.07, width - 0.2, 0.2, size=11, color=text_color, bold=True)


def add_bullets(slide, bullets, left, top, width, height, size=17, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.bullet = True
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_circle(slide, text, left, top, diameter, fill_color, text_color=WHITE, size=18):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    add_text_block(slide, text, left, top + diameter / 2 - 0.13, diameter, 0.3, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, left, top, width, height=0.32, fill_color=BLUE):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill_color
    arrow.line.fill.background()
    return arrow


def add_metric_card(slide, left, top, width, height, value, label, color):
    add_panel(slide, left, top, width, height, fill_color=WHITE)
    add_text_block(slide, value, left + 0.16, top + 0.18, width - 0.3, 0.45, size=24, color=color, bold=True, name="Aptos Display")
    add_text_block(slide, label, left + 0.16, top + 0.72, width - 0.3, height - 0.8, size=12, color=MUTED)


def add_family_card(slide, left, top, width, height, title, tag, body, color):
    add_panel(slide, left, top, width, height, fill_color=WHITE)
    add_badge(slide, tag, left + 0.14, top + 0.14, 1.05, fill_color=color, text_color=NAVY)
    add_text_block(slide, title, left + 0.14, top + 0.62, width - 0.28, 0.34, size=15, color=NAVY, bold=True, name="Aptos Display")
    add_text_block(slide, body, left + 0.14, top + 1.0, width - 0.28, height - 1.12, size=12, color=TEXT)


def add_timeline_node(slide, center_x, title, year, body, color):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(center_x - 0.18), Inches(3.08), Inches(0.36), Inches(0.36))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_block(slide, year, center_x - 0.45, 2.45, 0.9, 0.25, size=10, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_panel(slide, center_x - 1.0, 3.55, 2.0, 1.45, fill_color=WHITE)
    add_text_block(slide, title, center_x - 0.84, 3.72, 1.68, 0.3, size=14, color=NAVY, bold=True, name="Aptos Display", align=PP_ALIGN.CENTER)
    add_text_block(slide, body, center_x - 0.84, 4.08, 1.68, 0.75, size=10, color=TEXT, align=PP_ALIGN.CENTER)


def add_application_tile(slide, left, top, width, height, title, desc, color):
    add_panel(slide, left, top, width, height, fill_color=WHITE)
    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.18), Inches(top + 0.18), Inches(0.46), Inches(0.46))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    add_text_block(slide, title, left + 0.75, top + 0.16, width - 0.95, 0.3, size=14, color=NAVY, bold=True, name="Aptos Display")
    add_text_block(slide, desc, left + 0.18, top + 0.74, width - 0.36, height - 0.88, size=11, color=TEXT)


def add_comparison_bar(slide, left, top, label, pixel_value, latent_value):
    add_text_block(slide, label, left, top, 2.2, 0.22, size=11, color=MUTED)
    pixel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top + 0.32), Inches(pixel_value), Inches(0.18))
    pixel.fill.solid()
    pixel.fill.fore_color.rgb = GRAY_3
    pixel.line.fill.background()
    latent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top + 0.56), Inches(latent_value), Inches(0.18))
    latent.fill.solid()
    latent.fill.fore_color.rgb = TEAL
    latent.line.fill.background()


def add_reference_list(slide, refs, left, top, width):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.9))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, ref in enumerate(refs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ref
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_badge(slide, "Exposé scientifique", 0.86, 1.08, 1.9)
    add_text_block(slide, "Réseaux de diffusion", 0.86, 1.62, 5.7, 0.6, size=30, color=NAVY, bold=True, name="Aptos Display")
    add_text_block(
        slide,
        "Comprendre comment on génère une image en apprenant à enlever du bruit, puis pourquoi la diffusion latente a rendu ces modèles réellement pratiques.",
        0.88,
        2.26,
        5.5,
        1.05,
        size=16,
        color=TEXT,
    )

    add_panel(slide, 0.86, 3.58, 5.55, 2.05, fill_color=WHITE)
    add_text_block(slide, "Plan visuel", 1.06, 3.76, 1.7, 0.22, size=15, color=TEAL, bold=True, name="Aptos Display")
    add_bullets(
        slide,
        [
            "contexte et intuition",
            "problématique",
            "familles de solutions",
            "zoom sur Latent Diffusion Models",
            "résultats, limites, avenir",
        ],
        1.02,
        4.15,
        5.0,
        1.2,
        size=13,
    )

    add_panel(slide, 6.8, 1.1, 5.7, 4.95, fill_color=WHITE)
    add_text_block(slide, "Idée centrale", 7.06, 1.34, 2.0, 0.25, size=15, color=TEAL, bold=True, name="Aptos Display")
    stages = [
        ("Image", BLUE),
        ("+ bruit", GRAY_2),
        ("++ bruit", GRAY_3),
        ("bruit", GRAY_4),
        ("débruitage", TEAL),
    ]
    x = 7.0
    for index, (label, color) in enumerate(stages):
        add_circle(slide, "", x, 2.18, 0.78, color)
        add_text_block(slide, label, x - 0.12, 3.04, 1.05, 0.25, size=10, color=TEXT, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_arrow(slide, x + 0.84, 2.45, 0.55, fill_color=ORANGE if index >= 3 else BLUE)
        x += 1.42
    add_text_block(slide, "On va de la donnée vers le bruit, puis on apprend le chemin inverse.", 7.0, 3.7, 5.0, 0.5, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    add_text_block(slide, "Méthode détaillée: Rombach et al., CVPR 2022", 7.26, 4.55, 4.7, 0.25, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Support visuel original inspiré de ressources pédagogiques et appuyé sur des articles académiques")


def create_context_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "1. Partir d'une intuition connue", "Les blogs d'introduction insistent tous sur la même analogie: diffusion physique et débruitage progressif")
    add_panel(slide, 0.8, 1.95, 3.9, 4.35, fill_color=WHITE)
    add_text_block(slide, "Diffusion dans la vie courante", 1.02, 2.15, 2.6, 0.25, size=15, color=TEAL, bold=True, name="Aptos Display")
    add_circle(slide, "", 1.2, 2.85, 0.55, ORANGE)
    add_circle(slide, "", 1.95, 2.65, 0.48, GOLD)
    add_circle(slide, "", 2.65, 3.05, 0.4, MINT, text_color=TEXT)
    add_circle(slide, "", 3.15, 2.7, 0.32, GRAY_2, text_color=TEXT)
    add_text_block(slide, "Comme un parfum qui se répand dans une pièce: on passe d'une structure localisée à un état plus uniforme.", 1.04, 4.15, 3.3, 1.05, size=13)

    add_panel(slide, 4.95, 1.95, 3.55, 4.35, fill_color=WHITE)
    add_text_block(slide, "Version IA", 5.18, 2.15, 1.5, 0.25, size=15, color=TEAL, bold=True, name="Aptos Display")
    tones = [BLUE, RGBColor(113, 146, 199), GRAY_2, GRAY_3]
    x = 5.26
    for index, color in enumerate(tones):
        tile = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.75), Inches(0.63), Inches(0.63))
        tile.fill.solid()
        tile.fill.fore_color.rgb = color
        tile.line.fill.background()
        if index < len(tones) - 1:
            add_arrow(slide, x + 0.72, 2.95, 0.35)
        x += 1.0
    add_text_block(slide, "Image nette -> image bruitée -> image presque gaussienne.", 5.18, 4.15, 3.0, 0.45, size=13)
    add_text_block(slide, "Idée: apprendre ensuite le trajet inverse.", 5.18, 4.72, 2.8, 0.4, size=16, color=NAVY, bold=True)

    add_panel(slide, 8.75, 1.95, 3.75, 4.35, fill_color=MINT, line_color=MINT)
    add_text_block(slide, "Message à faire passer", 9.0, 2.16, 2.2, 0.25, size=15, color=TEAL, bold=True, name="Aptos Display")
    add_bullets(
        slide,
        [
            "Générer d'un seul coup est difficile.",
            "Débruiter étape par étape est plus simple.",
            "La diffusion transforme la génération en une suite de petits problèmes.",
        ],
        9.0,
        2.72,
        3.1,
        2.0,
        size=14,
    )
    add_footer(slide, "Inspiration pédagogique: Encord et Lilian Weng; formulation scientifique: Ho et al. 2020")


def create_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "2. Problématique", "Pourquoi la génération d'images réalistes reste un problème scientifique difficile")
    add_metric_card(slide, 0.9, 1.9, 2.45, 1.35, "786 432", "valeurs pour une image RGB 512 x 512", BLUE)
    add_metric_card(slide, 3.55, 1.9, 2.45, 1.35, "multi-mode", "une infinité de scènes plausibles pour un même prompt", ORANGE)
    add_metric_card(slide, 6.2, 1.9, 2.45, 1.35, "3 objectifs", "fidélité, diversité, contrôle", TEAL)
    add_metric_card(slide, 8.85, 1.9, 2.45, 1.35, "coût élevé", "entraînement et échantillonnage historiquement lents", GOLD)

    add_panel(slide, 1.0, 3.72, 5.4, 2.2, title="Le compromis central")
    axis_x = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.45), Inches(5.25), Inches(4.2), Inches(0.04))
    axis_x.fill.solid(); axis_x.fill.fore_color.rgb = MUTED; axis_x.line.fill.background()
    axis_y = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.45), Inches(4.25), Inches(0.04), Inches(1.04))
    axis_y.fill.solid(); axis_y.fill.fore_color.rgb = MUTED; axis_y.line.fill.background()
    add_text_block(slide, "qualité", 1.1, 4.0, 0.6, 0.2, size=10, color=MUTED)
    add_text_block(slide, "coût", 5.1, 5.35, 0.5, 0.2, size=10, color=MUTED)
    add_circle(slide, "GAN", 2.15, 4.95, 0.48, ORANGE, size=11)
    add_circle(slide, "VAE", 4.15, 5.02, 0.48, GRAY_3, size=11)
    add_circle(slide, "Diff.", 3.2, 4.45, 0.56, TEAL, size=11)
    add_text_block(slide, "Les diffusions se placent plus haut en qualité, mais avec un coût temporel important.", 1.7, 5.55, 4.1, 0.3, size=12)

    add_panel(slide, 6.75, 3.72, 5.55, 2.2, title="Question de recherche")
    add_text_block(slide, "Comment obtenir des images réalistes, diverses et contrôlables sans l'instabilité des GAN ni un coût prohibitif?", 6.98, 4.18, 5.0, 0.8, size=18, color=NAVY, bold=True, name="Aptos Display", align=PP_ALIGN.CENTER)
    add_footer(slide, "La suite de l'exposé répond à cette question par une revue de littérature puis un zoom sur LDM")


def create_mechanism_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "3. Mécanisme de base", "Schéma minimal d'un modèle de diffusion de type DDPM")
    add_panel(slide, 0.78, 1.92, 5.9, 4.6, title="Processus avant: ajouter du bruit")
    xs = [1.05, 2.15, 3.25, 4.35, 5.45]
    labels = ["x0", "x1", "x2", "...", "xT"]
    colors = [BLUE, RGBColor(105, 145, 214), GRAY_2, GRAY_3, GRAY_4]
    for idx, x in enumerate(xs):
        add_circle(slide, labels[idx], x, 3.0, 0.62, colors[idx], size=15 if labels[idx] != "..." else 14)
        if idx < len(xs) - 1:
            add_arrow(slide, x + 0.72, 3.17, 0.3)
    add_text_block(slide, "On applique T petites perturbations gaussiennes jusqu'à obtenir un bruit presque normal.", 1.06, 4.25, 5.1, 0.65, size=14)
    add_text_block(slide, "q(xt | xt-1)", 2.34, 2.55, 0.9, 0.2, size=10, color=MUTED)
    add_text_block(slide, "q(xt | xt-1)", 3.44, 2.55, 0.9, 0.2, size=10, color=MUTED)
    add_text_block(slide, "q(xt | xt-1)", 4.54, 2.55, 0.9, 0.2, size=10, color=MUTED)

    add_panel(slide, 6.86, 1.92, 5.68, 4.6, title="Processus inverse: enlever du bruit")
    xs = [11.6, 10.5, 9.4, 8.3, 7.2]
    labels = ["xT", "xT-1", "xT-2", "...", "x0"]
    colors = [GRAY_4, GRAY_3, GRAY_2, RGBColor(120, 177, 174), TEAL]
    for idx, x in enumerate(xs):
        add_circle(slide, labels[idx], x, 3.0, 0.62, colors[idx], size=11 if "-" in labels[idx] else 15)
        if idx < len(xs) - 1:
            add_arrow(slide, x - 0.48, 3.17, 0.3, fill_color=ORANGE)
    add_text_block(slide, "Le réseau prédit le bruit à retirer à chaque étape. Générer = parcourir le trajet inverse.", 7.16, 4.25, 4.9, 0.65, size=14)
    add_text_block(slide, "pθ(xt-1 | xt)", 10.23, 2.55, 0.95, 0.2, size=10, color=MUTED)
    add_text_block(slide, "pθ(xt-1 | xt)", 9.13, 2.55, 0.95, 0.2, size=10, color=MUTED)
    add_text_block(slide, "pθ(xt-1 | xt)", 8.03, 2.55, 0.95, 0.2, size=10, color=MUTED)
    add_footer(slide, "Inspiration visuelle: schémas pédagogiques vus chez Lilian Weng; contenu scientifique: Ho et al. 2020")


def create_family_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "4. Revue de littérature", "Situer les diffusions parmi les grandes familles de modèles génératifs")
    add_family_card(slide, 0.82, 1.95, 2.35, 2.0, "VAE", "stable", "Compression probabiliste et vraisemblance explicite, mais sorties souvent plus floues.", MINT)
    add_family_card(slide, 3.42, 1.95, 2.35, 2.0, "GAN", "net", "Très haute qualité visuelle, mais entraînement instable et risque de mode collapse.", ROSE)
    add_family_card(slide, 6.02, 1.95, 2.35, 2.0, "Flows", "exact", "Densité explicite et inversion exacte, au prix de contraintes architecturales.", RGBColor(230, 238, 251))
    add_family_card(slide, 8.62, 1.95, 2.35, 2.0, "Autoregressifs", "séquentiel", "Bonne vraisemblance, mais génération lente car un élément après l'autre.", RGBColor(249, 238, 219))
    add_family_card(slide, 11.22, 1.95, 1.3, 2.0, "Diff.", "robuste", "Qualité élevée et entraînement stable, mais coût de sampling.", MINT)

    add_panel(slide, 0.92, 4.35, 11.6, 1.7, title="Lecture rapide pour l'oral")
    add_text_block(slide, "Pourquoi les diffusions ont gagné en pratique?", 1.15, 4.7, 3.8, 0.25, size=15, color=NAVY, bold=True, name="Aptos Display")
    add_bullets(
        slide,
        [
            "elles évitent le duel générateur/discriminateur des GAN",
            "elles se prêtent bien au conditionnement par texte, classe ou masque",
            "leur qualité d'image a rapidement dépassé beaucoup d'alternatives",
        ],
        1.12,
        5.06,
        5.3,
        0.8,
        size=12,
    )
    add_text_block(slide, "Références de contexte", 7.05, 4.7, 2.0, 0.22, size=15, color=TEAL, bold=True, name="Aptos Display")
    add_text_block(slide, "Kingma and Welling 2014 | Goodfellow et al. 2014 | Dinh et al. 2018 | Ho et al. 2020", 7.05, 5.08, 4.8, 0.55, size=12, color=TEXT)
    add_footer(slide, "Cette slide remplit la partie revue de littérature à haut niveau")


def create_timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "5. Évolution récente des diffusions", "Une frise pour montrer que la méthode étudiée s'inscrit dans une dynamique très récente")
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(3.22), Inches(11.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LINE
    bar.line.fill.background()
    add_timeline_node(slide, 2.0, "DDPM", "2020", "Formalisation moderne\nFID 3.17 sur CIFAR-10", BLUE)
    add_timeline_node(slide, 4.3, "Score SDE", "2021", "Cadre continu\nFID 2.20", TEAL)
    add_timeline_node(slide, 6.6, "CFG", "2022", "Guidage conditionnel\nsans classifieur", ORANGE)
    add_timeline_node(slide, 8.9, "LDM", "2022", "Diffusion dans\nl'espace latent", GOLD)
    add_timeline_node(slide, 11.2, "Flow Matching", "2023", "Autres trajectoires\net ODE rapides", BLUE)
    add_footer(slide, "Références: Ho et al. 2020; Song et al. 2021; Ho and Salimans 2022; Rombach et al. 2022; Lipman et al. 2023")


def create_ldm_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "6. Méthode retenue: Latent Diffusion Models", "Méthode publiée à CVPR 2022, non marginale et à fort impact pratique")
    add_panel(slide, 0.82, 1.95, 11.8, 4.75, title="Pipeline visuel de la méthode")

    add_panel(slide, 1.08, 2.55, 2.05, 2.45, fill_color=WHITE)
    add_text_block(slide, "Image", 1.72, 2.82, 0.8, 0.25, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    for row in range(3):
        for col in range(3):
            tile = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.45 + col * 0.32), Inches(3.3 + row * 0.32), Inches(0.25), Inches(0.25))
            tile.fill.solid()
            tile.fill.fore_color.rgb = BLUE if (row + col) % 2 == 0 else RGBColor(114, 169, 216)
            tile.line.fill.background()

    add_arrow(slide, 3.3, 3.42, 0.55)

    add_panel(slide, 3.95, 2.35, 2.35, 2.85, fill_color=MINT, line_color=MINT)
    add_text_block(slide, "Autoencodeur", 4.46, 2.68, 1.32, 0.25, size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    add_text_block(slide, "Encode -> compresse\nDecode -> reconstruit", 4.34, 3.25, 1.6, 0.65, size=13, color=TEXT, align=PP_ALIGN.CENTER)
    add_text_block(slide, "But: retirer la redondance pixel", 4.2, 4.32, 1.88, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    add_arrow(slide, 6.45, 3.42, 0.55, fill_color=ORANGE)

    add_panel(slide, 7.12, 2.35, 2.45, 2.85, fill_color=WHITE)
    add_text_block(slide, "Latent z", 7.86, 2.68, 1.0, 0.25, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    for row in range(2):
        for col in range(4):
            tile = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.5 + col * 0.34), Inches(3.35 + row * 0.34), Inches(0.22), Inches(0.22))
            tile.fill.solid()
            tile.fill.fore_color.rgb = TEAL if (row + col) % 2 == 0 else RGBColor(135, 208, 202)
            tile.line.fill.background()
    add_text_block(slide, "Espace compact mais sémantiquement riche", 7.34, 4.3, 2.0, 0.45, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    add_arrow(slide, 9.72, 3.42, 0.55)

    add_panel(slide, 10.4, 2.35, 1.95, 2.85, fill_color=ROSE, line_color=ROSE)
    add_text_block(slide, "U-Net\n+ cross-attn", 10.74, 2.85, 1.26, 0.6, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    add_text_block(slide, "Débruite dans z\nConditionne par texte", 10.68, 3.82, 1.35, 0.6, size=12, color=TEXT, align=PP_ALIGN.CENTER)

    add_text_block(slide, "Famille: diffusion conditionnelle. Rupture: diffusion en latent plutôt qu'en pixels.", 1.05, 5.8, 11.2, 0.45, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER, name="Aptos Display")
    add_footer(slide, "Rombach et al. High-Resolution Image Synthesis with Latent Diffusion Models. CVPR 2022")


def create_efficiency_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "7. Pourquoi cette méthode a marqué le domaine", "Le gain clé n'est pas seulement la qualité, c'est le compromis qualité / coût / contrôle")
    add_panel(slide, 0.88, 1.92, 5.2, 4.7, title="Pixel-space vs latent-space")
    add_text_block(slide, "Pixel-space diffusion", 1.16, 2.42, 1.8, 0.24, size=14, color=MUTED, bold=True)
    add_text_block(slide, "Latent diffusion", 3.86, 2.42, 1.5, 0.24, size=14, color=TEAL, bold=True)
    add_comparison_bar(slide, 1.16, 2.9, "mémoire", 3.2, 1.45)
    add_comparison_bar(slide, 1.16, 3.72, "temps d'entraînement", 3.8, 1.7)
    add_comparison_bar(slide, 1.16, 4.54, "résolution praticable", 2.25, 3.2)
    add_text_block(slide, "Les auteurs expliquent que les diffusions pixel-space demandent souvent des centaines de GPU-days à haute résolution.", 1.16, 5.42, 4.45, 0.72, size=12)

    add_panel(slide, 6.35, 1.92, 6.0, 4.7, title="Ce que l'on gagne")
    add_bullets(
        slide,
        [
            "compression perceptuelle avant modélisation générative",
            "cross-attention pour texte, cartes sémantiques, boîtes",
            "qualité visuelle forte avec coût bien plus réaliste",
            "base conceptuelle de Stable Diffusion",
        ],
        6.62,
        2.45,
        5.2,
        1.8,
        size=15,
    )
    add_text_block(slide, "Résultat pratique", 6.66, 4.82, 2.0, 0.24, size=15, color=TEAL, bold=True, name="Aptos Display")
    add_text_block(slide, "La diffusion devient exploitable pour le text-to-image grand public sans abandonner la qualité.", 6.66, 5.18, 5.2, 0.52, size=18, color=NAVY, bold=True, name="Aptos Display")
    add_footer(slide, "Inspiration pédagogique: billet Medium sur LDM; justification scientifique: résumé et expériences du papier CVPR 2022")


def create_applications_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "8. Applications et descendants", "Slide volontairement visuelle pour alléger la densité de texte")
    add_application_tile(slide, 0.9, 1.95, 3.75, 1.75, "Text-to-image", "DALL-E 2, Imagen, Stable Diffusion: le texte devient une condition du débruitage.", BLUE)
    add_application_tile(slide, 4.78, 1.95, 3.75, 1.75, "Inpainting", "On masque une zone puis on reconstruit de façon cohérente avec le reste de l'image.", TEAL)
    add_application_tile(slide, 8.66, 1.95, 3.75, 1.75, "Super-résolution", "On part d'une image basse résolution et on ajoute du détail plausible.", ORANGE)
    add_application_tile(slide, 0.9, 3.95, 3.75, 1.75, "Image-to-image", "Colorisation, stylisation, transformation guidée par une image source.", GOLD)
    add_application_tile(slide, 4.78, 3.95, 3.75, 1.75, "Vidéo", "Le défi devient la cohérence temporelle d'une frame à la suivante.", BLUE)
    add_application_tile(slide, 8.66, 3.95, 3.75, 1.75, "Contrôle structuré", "ControlNet et cartes de profondeur ou contours pour guider la génération.", TEAL)
    add_footer(slide, "Inspiration: panorama d'applications présenté par Encord et Lilian Weng; références scientifiques citées oralement")


def create_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "9. Résultats et lecture critique", "Comment présenter les résultats sans noyer l'auditoire dans les métriques")
    add_panel(slide, 0.9, 1.95, 3.85, 4.7, title="Résultats à retenir")
    add_bullets(
        slide,
        [
            "DDPM 2020: FID 3.17 sur CIFAR-10",
            "Score SDE 2021: FID 2.20 sur CIFAR-10",
            "LDM 2022: résultats très compétitifs sur inpainting, synthèse sémantique et super-résolution",
        ],
        1.16,
        2.45,
        3.2,
        2.1,
        size=14,
    )
    add_text_block(slide, "Message oral: les chiffres montrent la progression, mais l'impact réel de LDM est surtout computationnel et applicatif.", 1.15, 5.25, 3.2, 0.8, size=12)

    add_panel(slide, 5.0, 1.95, 3.45, 4.7, title="Forces")
    add_bullets(
        slide,
        [
            "très bonne qualité visuelle",
            "entraînement plus stable que GAN",
            "conditionnement flexible",
            "forte modularité",
        ],
        5.22,
        2.46,
        2.8,
        1.8,
        size=14,
    )
    add_panel(slide, 8.7, 1.95, 3.75, 4.7, title="Faiblesses et avenir")
    add_bullets(
        slide,
        [
            "sampling encore lent",
            "biais hérités des données",
            "enjeux de gouvernance et d'usage",
            "travaux récents: consistency models, DiT, flow matching",
        ],
        8.95,
        2.46,
        3.0,
        2.1,
        size=14,
    )
    add_footer(slide, "Cette slide couvre la conclusion scientifique demandée: forces, faiblesses et travaux à venir")


def create_references_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "10. Références", "Séparer clairement les sources académiques des sources d'inspiration pédagogique")
    add_panel(slide, 0.88, 1.92, 7.55, 4.95, title="Articles académiques à citer")
    add_reference_list(
        slide,
        [
            "Sohl-Dickstein et al. Deep Unsupervised Learning using Nonequilibrium Thermodynamics. ICML 2015.",
            "Ho, Jain, Abbeel. Denoising Diffusion Probabilistic Models. NeurIPS 2020.",
            "Song et al. Score-Based Generative Modeling through Stochastic Differential Equations. ICLR 2021.",
            "Ho, Salimans. Classifier-Free Diffusion Guidance. 2022.",
            "Rombach et al. High-Resolution Image Synthesis with Latent Diffusion Models. CVPR 2022.",
            "Lipman et al. Flow Matching for Generative Modeling. ICLR 2023.",
            "Peebles and Xie. Scalable Diffusion Models with Transformers. ICCV 2023.",
        ],
        1.1,
        2.35,
        6.95,
    )

    add_panel(slide, 8.7, 1.92, 3.72, 4.95, title="Sources d'inspiration")
    add_bullets(
        slide,
        [
            "Encord blog: bonne entrée en matière pour l'intuition et les applications",
            "Lilian Weng: excellente synthèse technique et historique",
            "Medium sur LDM: utile pour vulgariser le pipeline latent",
            "LinkedIn fourni mais contenu non récupérable sans connexion",
        ],
        8.95,
        2.38,
        3.0,
        2.2,
        size=13,
    )
    add_text_block(slide, "Important: ces pages servent d'inspiration visuelle et pédagogique, mais les citations de l'exposé doivent rester académiques.", 8.95, 5.35, 2.95, 0.85, size=12, color=NAVY, bold=True)
    add_footer(slide, "Aucune image externe n'est recopiée; le support utilise des schémas originaux")


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs)
    create_context_slide(prs)
    create_problem_slide(prs)
    create_mechanism_slide(prs)
    create_family_slide(prs)
    create_timeline_slide(prs)
    create_ldm_slide(prs)
    create_efficiency_slide(prs)
    create_applications_slide(prs)
    create_results_slide(prs)
    create_references_slide(prs)

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation created: {OUTPUT_PATH}")
